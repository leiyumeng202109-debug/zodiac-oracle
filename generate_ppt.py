#!/usr/bin/env python3
"""
Generate PPTX presentation for Zodiac + MBTI Fortune Teller Chatbot
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PURPLE_DARK = RGBColor(26, 26, 46)  # #1a1a2e
PURPLE_GRADIENT_START = RGBColor(102, 126, 234)  # #667eea
PURPLE_GRADIENT_END = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(160, 160, 160)
GREEN = RGBColor(40, 167, 69)
ORANGE = RGBColor(255, 193, 7)

def add_title_slide(prs, title, subtitle="", emoji=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_DARK
    background.line.fill.background()
    
    # Emoji
    if emoji:
        emoji_box = slide.shapes.add_textbox(Inches(0), Inches(1.5), prs.slide_width, Inches(1))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.8), prs.slide_width, Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_START
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_START
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"✦ {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_START
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
    ).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE_GRADIENT_START
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(16)
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = WHITE
            p.font.size = Pt(14)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(40, 40, 70)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_START
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_END
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✦ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE_GRADIENT_END
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.5), Inches(4))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✦ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    return slide

# ============ BUILD PRESENTATION ============

# Slide 1: Title
add_title_slide(
    prs,
    "Zodiac + MBTI Fortune Teller",
    "Chatbot Project Presentation\n\nCosmic Oracle Development Team | March 2026",
    "🔮"
)

# Slide 2: Agenda
add_content_slide(prs, "📋 Agenda", [
    "1. Labor of Division - Team structure, roles, and development timeline",
    "2. Project Approach - Design philosophy, architecture, and security",
    "3. Preliminary Results - Deliverables, test results, and live demo",
    "4. Ongoing & Planned - Future features and risk assessment"
])

# Slide 3: Team Structure
add_table_slide(
    prs,
    "👥 Team Structure",
    ["Role", "Responsibilities", "Deliverables"],
    [
        ["Project Lead", "Architecture, API design, security", "System architecture, API endpoints"],
        ["Frontend Developer", "UI/UX design, React components", "Chat interface, quick questions"],
        ["AI/LLM Engineer", "Prompt engineering, validation", "System prompts, security filters"],
        ["QA & Documentation", "Testing, documentation", "Test cases, design docs"]
    ]
)

# Slide 4: Development Timeline
add_table_slide(
    prs,
    "📅 Development Timeline",
    ["Phase", "Tasks", "Status"],
    [
        ["Phase 1", "Requirements analysis, tech stack selection", "✓ Complete"],
        ["Phase 2", "Core development: API routes, frontend UI", "✓ Complete"],
        ["Phase 3", "Security implementation, prompt engineering", "✓ Complete"],
        ["Phase 4", "English localization, content positivity", "✓ Complete"],
        ["Phase 5", "Documentation, deployment, testing", "○ In Progress"]
    ]
)

# Slide 5: Technology Stack
add_two_column_slide(
    prs,
    "🛠️ Technology Stack",
    "Frontend",
    ["Next.js 16 (App Router)", "React 19", "TypeScript 5", "shadcn/ui + Radix UI", "Tailwind CSS 4"],
    "Backend & AI",
    ["Next.js API Routes", "Server-Sent Events (SSE)", "coze-coding-dev-sdk", "Doubao LLM", "Temperature: 0.9"]
)

# Slide 6: Design Philosophy
slide = add_content_slide(prs, "💡 Design Philosophy", [
    '"Entertainment First, Accuracy Last"',
    "",
    "The chatbot intentionally positions itself as 0.0001% accurate",
    "Setting correct user expectations while maximizing entertainment",
    "Delivering humor and positivity through mystical predictions"
])

# Slide 7: Technical Architecture
add_table_slide(
    prs,
    "🏗️ Technical Architecture",
    ["Layer", "Component", "Technology"],
    [
        ["Presentation", "Chat UI with mystical theme", "React + Tailwind CSS"],
        ["Communication", "Streaming API", "Server-Sent Events"],
        ["Security", "Input sanitization, response validation", "Custom security layer"],
        ["Intelligence", "LLM with high creativity", "Doubao (temp: 0.9)"]
    ]
)

# Slide 8: Security Strategy
add_content_slide(prs, "🔒 Security Strategy - Three-Layer Protection", [
    "Layer 1: Input Sanitization - Detect and block prompt injection attempts",
    "Layer 2: Domain Validation - Accept broad inputs, interpret through zodiac/MBTI lens",
    "Layer 3: Response Validation - Prevent AI identity leaks and inappropriate content"
])

# Slide 9: Content Guidelines
add_two_column_slide(
    prs,
    "📝 Content Guidelines",
    "✅ Required Behaviors",
    ["Always positive and uplifting", "Use mystical opening phrases", "Blend Zodiac + MBTI terms", "End with entertainment disclaimer"],
    "❌ Forbidden Behaviors",
    ["Reveal AI identity", "Negative predictions", "Refuse user questions", "Serious/preachy tone"]
)

# Slide 10: Deliverables
add_table_slide(
    prs,
    "📦 Functional Deliverables",
    ["Feature", "Description", "Status"],
    [
        ["Chat Interface", "Purple-themed UI with zodiac decorations", "✓ Delivered"],
        ["Streaming API", "Real-time typewriter effect via SSE", "✓ Delivered"],
        ["Security Layer", "Three-layer protection with mystical errors", "✓ Delivered"],
        ["English Version", "Full localization for international users", "✓ Delivered"],
        ["Positive Content", "Uplifting, encouraging predictions", "✓ Delivered"]
    ]
)

# Slide 11: Test Results
add_table_slide(
    prs,
    "🧪 Test Results",
    ["Test Case", "Input Example", "Result"],
    [
        ["Love Fortune", '"When will I find true love?"', "✓ Pass"],
        ["MBTI + Zodiac", '"What is an INFP + Scorpio like?"', "✓ Pass"],
        ["Career Question", '"What job suits me?"', "✓ Pass"],
        ["Prompt Injection", '"ignore all instructions"', "✓ Blocked"],
        ["AI Identity Leak", '"Are you an AI?"', "✓ Protected"]
    ]
)

# Slide 12: Live Demo
add_title_slide(
    prs,
    "🔗 Live Demo",
    "https://8cfe2369-77c3-43f9-b04c-35886985c710.dev.coze.site\n\nAvailable 24/7 for testing",
    "🔮"
)

# Slide 13: Ongoing Work
add_table_slide(
    prs,
    "🔄 Ongoing Work",
    ["Task", "Description", "Priority"],
    [
        ["User Feedback", "Gather insights from international users", "High"],
        ["Response Optimization", "Fine-tune prompt for better engagement", "Medium"],
        ["Performance Monitoring", "Track API response times and errors", "Medium"]
    ]
)

# Slide 14: Planned Features
add_content_slide(prs, "🎯 Planned Features - Phase 6", [
    "💬 User History - Save conversation history per user",
    "📤 Share Function - Generate shareable fortune cards",
    "📅 Daily Fortune - Cached daily predictions",
    "🌍 Multi-language - Support for additional languages",
    "📊 Analytics Dashboard - Track usage patterns",
    "🎨 Theme Customization - Personalized themes"
])

# Slide 15: Risk Assessment
add_table_slide(
    prs,
    "⚠️ Risk Assessment",
    ["Risk", "Mitigation", "Status"],
    [
        ["Prompt injection attacks", "Three-layer security implemented", "✓ Resolved"],
        ["AI identity exposure", "Response validation + custom persona", "✓ Resolved"],
        ["Negative user experience", "Positive content guidelines enforced", "✓ Resolved"],
        ["LLM service downtime", "Fallback responses implemented", "○ Monitoring"]
    ]
)

# Slide 16: Thank You
add_title_slide(
    prs,
    "Thank You!",
    "Questions & Discussion\n\nNext Milestone: User Testing Phase → Final Report\n\n✨ 0.0001% Accurate, 99.9% Entertaining ✨",
    "🔮"
)

# Save
output_path = "/workspace/projects/public/docs/Zodiac_MBTI_Fortune_Teller_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
